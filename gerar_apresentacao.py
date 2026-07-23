"""
Gera a apresentação PPT do sistema Gestão de Vínculo — SenseBike
Execute: python gerar_apresentacao.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Paleta SenseBike ────────────────────────────────────────────────────────
VERDE       = RGBColor(0x2E, 0x7D, 0x32)   # verde escuro
VERDE_CLARO = RGBColor(0x43, 0xA0, 0x47)   # verde médio
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESC   = RGBColor(0x21, 0x21, 0x21)
CINZA_MED   = RGBColor(0x55, 0x55, 0x55)
CINZA_CLR   = RGBColor(0xF5, 0xF5, 0xF5)
AMARELO     = RGBColor(0xFF, 0xC1, 0x07)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=CINZA_ESC,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def slide_capa(prs):
    sl = blank_slide(prs)
    bg(sl, VERDE)

    # Faixa lateral esquerda clara
    box(sl, 0, 0, Inches(0.18), H, VERDE_CLARO)

    # Título
    txt(sl, "Gestão de Vínculo", Inches(0.6), Inches(1.6), Inches(9), Inches(1.4),
        size=52, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

    # Subtítulo
    txt(sl, "Sistema de Pedidos de Vínculo entre Franquias e Matriz",
        Inches(0.6), Inches(3.0), Inches(10), Inches(0.8),
        size=24, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.LEFT)

    # Data
    data = datetime.date.today().strftime("%d/%m/%Y")
    txt(sl, f"Apresentação Gerência  |  {data}",
        Inches(0.6), Inches(6.5), Inches(8), Inches(0.6),
        size=14, color=RGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.LEFT)

    # Logo texto canto direito
    txt(sl, "SenseBike", Inches(10.5), Inches(6.4), Inches(2.5), Inches(0.8),
        size=20, bold=True, color=AMARELO, align=PP_ALIGN.RIGHT)


def slide_problema(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "O Problema que Resolvemos",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

    items = [
        ("📋", "Pedidos de vínculo feitos por e-mail ou WhatsApp, sem rastreabilidade"),
        ("🔄", "Nenhuma visibilidade de em qual etapa o pedido estava parado"),
        ("📎", "Anexos e documentos dispersos, difíceis de localizar depois"),
        ("❌", "Sem registro formal de reprovações e justificativas"),
        ("⏱️", "Demora para acionamento do time correto (Financeiro / TI)"),
    ]

    for i, (icon, desc) in enumerate(items):
        top = Inches(1.45) + i * Inches(1.0)
        box(sl, Inches(0.4), top + Inches(0.05), Inches(0.55), Inches(0.7), VERDE_CLARO)
        txt(sl, icon, Inches(0.42), top, Inches(0.6), Inches(0.8), size=22, align=PP_ALIGN.CENTER)
        txt(sl, desc, Inches(1.1), top + Inches(0.08), Inches(11.5), Inches(0.7),
            size=18, color=CINZA_ESC)


def slide_solucao(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "A Solução: Plataforma Web Centralizada",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO)

    cards = [
        (VERDE,       "📊", "Dashboard\npor Perfil",   "Cada equipe vê\nsomente o que\nlhe compete"),
        (RGBColor(0x15,0x65,0xC0), "🔁", "Fluxo\nAutomatizado", "Status avança\nautomaticamente\na cada aprovação"),
        (RGBColor(0x6A,0x1B,0x9A), "📎", "Anexos na\nNuvem",     "Documentos salvos\nno Cloudinary,\nsempre acessíveis"),
        (RGBColor(0xE6,0x59,0x00), "📧", "Notificações\npor E-mail", "Time avisado\nautomaticamente\nno momento certo"),
    ]

    cw = Inches(2.9)
    for i, (color, icon, title, desc) in enumerate(cards):
        l = Inches(0.35) + i * Inches(3.25)
        box(sl, l, Inches(1.4), cw, Inches(4.8), color)
        txt(sl, icon,  l + Inches(0.2), Inches(1.55), cw, Inches(0.9), size=36, color=BRANCO)
        txt(sl, title, l + Inches(0.2), Inches(2.5),  cw, Inches(1.1), size=20, bold=True, color=BRANCO)
        txt(sl, desc,  l + Inches(0.2), Inches(3.65), cw, Inches(2.0), size=15, color=RGBColor(0xE0,0xE0,0xE0))


def slide_fluxo(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "Fluxo de Aprovação",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO)

    etapas = [
        (AMARELO,                          CINZA_ESC, "1", "Franquia\nenvia pedido"),
        (RGBColor(0x15,0x65,0xC0), BRANCO,            "2", "Comercial\nanalisa"),
        (RGBColor(0x6A,0x1B,0x9A), BRANCO,            "3", "Financeiro\nvalida (se necessário)"),
        (RGBColor(0x00,0x83,0x8A), BRANCO,            "4", "TI\nexecuta"),
        (VERDE,                    BRANCO,             "✓", "Pedido\nFechado"),
    ]

    bw = Inches(2.1)
    for i, (color, tcol, num, label) in enumerate(etapas):
        l = Inches(0.3) + i * Inches(2.6)
        box(sl, l, Inches(2.0), bw, Inches(2.4), color)
        txt(sl, num,   l, Inches(2.05), bw, Inches(0.9), size=38, bold=True, color=tcol, align=PP_ALIGN.CENTER)
        txt(sl, label, l, Inches(2.95), bw, Inches(1.2), size=16, bold=True, color=tcol, align=PP_ALIGN.CENTER)

        # seta →
        if i < len(etapas) - 1:
            txt(sl, "→", l + bw, Inches(2.5), Inches(0.5), Inches(0.8),
                size=28, color=CINZA_MED, align=PP_ALIGN.CENTER)

    # Nota reprovação
    box(sl, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.6), RGBColor(0xFF,0xEB,0xEE))
    txt(sl, "↩  Reprovação:",
        Inches(0.5), Inches(5.1), Inches(3), Inches(0.6),
        size=16, bold=True, color=RGBColor(0xC6,0x28,0x28))
    txt(sl, "Em qualquer etapa o pedido pode ser reprovado — volta ao status Aberto com "
            "justificativa registrada e a franquia (ou comercial) é notificada por e-mail.",
        Inches(0.5), Inches(5.65), Inches(12.5), Inches(0.8),
        size=15, color=CINZA_ESC)


def slide_perfis(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "Perfis de Acesso",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO)

    perfis = [
        ("🏪 Franquia",    AMARELO,                   CINZA_ESC, "Abre novos pedidos de vínculo\nAcompanha status dos seus pedidos"),
        ("💼 Comercial",   VERDE,                     BRANCO,    "Vê todos os pedidos\nAprova ou reprova (avança para Financeiro ou TI)"),
        ("💰 Financeiro",  RGBColor(0x15,0x65,0xC0),  BRANCO,    "Vê pedidos aguardando validação financeira\nAprova com possibilidade de anexar documentos"),
        ("🖥️ TI",          RGBColor(0x6A,0x1B,0x9A), BRANCO,    "Vê pedidos prontos para execução\nAprova ao concluir a tarefa técnica"),
        ("⚙️ Admin",       CINZA_ESC,                 BRANCO,    "Acesso total ao sistema\nGerencia configurações de SMTP e templates de e-mail"),
    ]

    rh = Inches(1.08)
    for i, (nome, color, tcol, desc) in enumerate(perfis):
        top = Inches(1.35) + i * rh
        box(sl, Inches(0.3), top, Inches(2.8), rh - Inches(0.06), color)
        txt(sl, nome, Inches(0.35), top + Inches(0.2), Inches(2.7), Inches(0.7),
            size=17, bold=True, color=tcol)

        box(sl, Inches(3.2), top, Inches(9.8), rh - Inches(0.06), RGBColor(0xEE,0xEE,0xEE))
        txt(sl, desc, Inches(3.35), top + Inches(0.12), Inches(9.5), Inches(0.88),
            size=15, color=CINZA_ESC)


def slide_notificacoes(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "Notificações Automáticas por E-mail",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO)

    eventos = [
        ("📨", "Novo pedido recebido",       "Comercial é notificado assim que a franquia abre o pedido"),
        ("💰", "Aprovado para Financeiro",    "Time financeiro recebe alerta para iniciar a análise"),
        ("🖥️", "Aprovado para TI",           "Time de TI é acionado para execução da tarefa"),
        ("✅", "Pedido vinculado (fechado)",  "Franquia recebe confirmação de que o vínculo foi concluído"),
        ("❌", "Pedido reprovado",            "Destinatário correto (franquia ou comercial) recebe a justificativa"),
    ]

    for i, (icon, titulo, desc) in enumerate(eventos):
        top = Inches(1.4) + i * Inches(1.1)
        box(sl, Inches(0.3), top, Inches(0.7), Inches(0.85), VERDE_CLARO)
        txt(sl, icon, Inches(0.3), top, Inches(0.7), Inches(0.85), size=22, align=PP_ALIGN.CENTER)
        txt(sl, titulo, Inches(1.15), top + Inches(0.02), Inches(4.0), Inches(0.45),
            size=16, bold=True, color=CINZA_ESC)
        txt(sl, desc, Inches(1.15), top + Inches(0.45), Inches(11.5), Inches(0.45),
            size=14, color=CINZA_MED)


def slide_tecnologia(prs):
    sl = blank_slide(prs)
    bg(sl, CINZA_CLR)
    box(sl, 0, 0, W, Inches(1.2), VERDE)

    txt(sl, "Tecnologia & Infraestrutura",
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        size=32, bold=True, color=BRANCO)

    cols = [
        ("Frontend", VERDE, [
            "Next.js 14 (React)",
            "TypeScript",
            "TailwindCSS",
            "React Query",
        ]),
        ("Backend", RGBColor(0x15,0x65,0xC0), [
            "FastAPI (Python)",
            "SQLAlchemy async",
            "JWT + refresh token",
            "Alembic migrations",
        ]),
        ("Dados", RGBColor(0x6A,0x1B,0x9A), [
            "PostgreSQL",
            "Redis (cache)",
            "Cloudinary (anexos)",
            "Docker Compose",
        ]),
        ("Operação", RGBColor(0x00,0x83,0x8A), [
            "Deploy via Coolify",
            "SMTP Gmail configurável",
            "Templates editáveis",
            "Logs centralizados",
        ]),
    ]

    cw = Inches(2.9)
    for i, (titulo, color, items) in enumerate(cols):
        l = Inches(0.35) + i * Inches(3.25)
        box(sl, l, Inches(1.4), cw, Inches(0.6), color)
        txt(sl, titulo, l, Inches(1.42), cw, Inches(0.56),
            size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

        box(sl, l, Inches(2.0), cw, Inches(4.0), RGBColor(0xEE,0xEE,0xEE))
        for j, item in enumerate(items):
            txt(sl, f"• {item}", l + Inches(0.15), Inches(2.1) + j * Inches(0.85), cw, Inches(0.8),
                size=16, color=CINZA_ESC)


def slide_encerramento(prs):
    sl = blank_slide(prs)
    bg(sl, VERDE)
    box(sl, 0, 0, Inches(0.18), H, VERDE_CLARO)

    txt(sl, "Sistema pronto para uso.",
        Inches(0.7), Inches(1.8), Inches(11), Inches(1.2),
        size=46, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

    txt(sl, "Fluxo digital  •  Rastreabilidade total  •  Equipes notificadas automaticamente",
        Inches(0.7), Inches(3.1), Inches(11.5), Inches(0.8),
        size=22, color=RGBColor(0xC8,0xE6,0xC9), align=PP_ALIGN.LEFT)

    txt(sl, "Obrigado!",
        Inches(0.7), Inches(5.0), Inches(8), Inches(1.0),
        size=34, bold=True, color=AMARELO, align=PP_ALIGN.LEFT)

    data = datetime.date.today().strftime("%d/%m/%Y")
    txt(sl, f"SenseBike  |  {data}",
        Inches(0.7), Inches(6.5), Inches(8), Inches(0.6),
        size=14, color=RGBColor(0xA5,0xD6,0xA7), align=PP_ALIGN.LEFT)


def main():
    prs = prs_new()
    slide_capa(prs)
    slide_problema(prs)
    slide_solucao(prs)
    slide_fluxo(prs)
    slide_perfis(prs)
    slide_notificacoes(prs)
    slide_tecnologia(prs)
    slide_encerramento(prs)

    out = "Gestao_de_Vinculo_SenseBike.pptx"
    prs.save(out)
    print(f"OK: Apresentacao gerada: {out}")
    print(f"   {prs.slides.__len__()} slides | 16:9 widescreen")


if __name__ == "__main__":
    main()
